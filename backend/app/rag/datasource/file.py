import logging
import docx
import pptx
import os
import openpyxl
from PIL import Image as PILImage
import io
from pydantic import BaseModel
from typing import Generator, IO
from pypdf import PdfReader

from app.rag.image_analysis import LLMImageAnalyzer

from app.models import Document, Upload
from app.models.image import Image
from app.file_storage import default_file_storage
from app.types import MimeTypes
from .base import BaseDataSource

logger = logging.getLogger(__name__)


class FileConfig(BaseModel):
    file_id: int


class FileDataSource(BaseDataSource):
    def validate_config(self):
        if not isinstance(self.config, list):
            raise ValueError("config must be a list")
        for f_config in self.config:
            FileConfig.model_validate(f_config)

    def load_documents(self) -> Generator[Document, None, None]:
        logger.info(f"Starting document loading for data source ID: {self.data_source_id}")
        for f_config in self.config:
            upload_id = f_config["file_id"]
            logger.info(f"Processing file with upload ID: {upload_id}")
            upload = self.session.get(Upload, upload_id)
            if upload is None:
                logger.error(f"Upload with id {upload_id} not found")
                continue
            logger.info(f"Found upload: {upload.name}, mime type: {upload.mime_type}")

            with default_file_storage.open(upload.path) as f:
                logger.debug(f"Opened file: {upload.path}")
                if upload.mime_type == MimeTypes.PDF:
                    logger.info(f"Processing PDF file: {upload.name}")
                    # Create document first to get an ID
                    document = Document(
                        name=upload.name,
                        hash=hash(f.read()),  # Generate hash from file content
                        content="",  # Will be populated after extraction
                        mime_type=MimeTypes.PLAIN_TXT,
                        data_source_id=self.data_source_id,
                        user_id=self.user_id,
                        source_uri=upload.path,
                        last_modified_at=upload.created_at,
                    )
                    logger.debug(f"Created document object for {upload.name}")
                    self.session.add(document)
                    self.session.flush()  # Get the document ID
                    logger.info(f"Assigned document ID: {document.id} for {upload.name}")
                    
                    # Reset file pointer
                    f.seek(0)
                    logger.debug(f"Reset file pointer for {upload.name}")
                    
                    # Extract text and images
                    try:
                        logger.info(f"Starting text and image extraction for document ID: {document.id}")
                        content, images = extract_text_and_images_from_pdf(f, document.id)
                        logger.info(f"Extraction completed for document ID: {document.id}, found {len(images)} images")
                        
                        # Update document with content
                        document.content = content
                        logger.debug(f"Updated document content, length: {len(content)}")
                        
                        # Save images to database
                        if images:
                            logger.info(f"Adding {len(images)} images to session for document ID: {document.id}")
                            self.session.add_all(images)
                            self.session.flush()
                            logger.info(f"Successfully added {len(images)} images to database for document ID: {document.id}")
                        else:
                            logger.warning(f"No images extracted from PDF document ID: {document.id}")
                    except Exception as e:
                        logger.error(f"Error during PDF extraction for document ID: {document.id}: {str(e)}")
                        logger.error("Stack trace:", exc_info=True)
                        # Still yield the document even if image extraction failed
                        document.content = "Error extracting content from PDF"
                    
                    yield document
                    logger.info(f"Yielded document ID: {document.id}")
                elif upload.mime_type == MimeTypes.DOCX:
                    content = extract_text_from_docx(f)
                    mime_type = MimeTypes.PLAIN_TXT
                    document = Document(
                        name=upload.name,
                        hash=hash(content),
                        content=content,
                        mime_type=mime_type,
                        data_source_id=self.data_source_id,
                        user_id=self.user_id,
                        source_uri=upload.path,
                        last_modified_at=upload.created_at,
                    )
                    yield document
                elif upload.mime_type == MimeTypes.PPTX:
                    content = extract_text_from_pptx(f)
                    mime_type = MimeTypes.PLAIN_TXT
                    document = Document(
                        name=upload.name,
                        hash=hash(content),
                        content=content,
                        mime_type=mime_type,
                        data_source_id=self.data_source_id,
                        user_id=self.user_id,
                        source_uri=upload.path,
                        last_modified_at=upload.created_at,
                    )
                    yield document
                elif upload.mime_type == MimeTypes.XLSX:
                    content = extract_text_from_xlsx(f)
                    mime_type = MimeTypes.PLAIN_TXT
                    document = Document(
                        name=upload.name,
                        hash=hash(content),
                        content=content,
                        mime_type=mime_type,
                        data_source_id=self.data_source_id,
                        user_id=self.user_id,
                        source_uri=upload.path,
                        last_modified_at=upload.created_at,
                    )
                    yield document
                else:
                    content = f.read()
                    mime_type = upload.mime_type
                    document = Document(
                        name=upload.name,
                        hash=hash(content),
                        content=content,
                        mime_type=mime_type,
                        data_source_id=self.data_source_id,
                        user_id=self.user_id,
                        source_uri=upload.path,
                        last_modified_at=upload.created_at,
                    )
                    yield document


def extract_text_from_pdf(file: IO) -> str:
    reader = PdfReader(file)
    print("\n\n".join([page.extract_text() for page in reader.pages]))
    return "\n\n".join([page.extract_text() for page in reader.pages])

def extract_text_and_images_from_pdf(file: IO, document_id: int) -> tuple[str, list[Image]]:
    reader = PdfReader(file)
    text_content = []
    images_list = []
    
    # Import here to avoid circular imports
    from pdf2image import convert_from_bytes
    import uuid
    import os
    from app.rag.chat_config import get_default_llm
    from sqlmodel import Session
    from app.core.db import engine
    
    logger.info(f"Starting PDF extraction for document ID: {document_id}")
    
    # Initialize LLM image analyzer
    try:
        with Session(engine) as session:
            llm = get_default_llm(session)
            image_analyzer = LLMImageAnalyzer(llm)
    except Exception as e:
        logger.error(f"Failed to initialize LLM image analyzer: {str(e)}")
        raise
    
    for page_num, page in enumerate(reader.pages):
        logger.info(f"Processing page {page_num + 1} for document ID: {document_id}")
        
        # Extract embedded text from the page
        try:
            embedded_text = page.extract_text().strip()
            logger.debug(f"Extracted embedded text from page {page_num + 1}, length: {len(embedded_text)}")
        except Exception as e:
            logger.error(f"Error extracting embedded text from page {page_num + 1}: {str(e)}")
            embedded_text = ""
        
        # Convert the entire PDF file to images
        file.seek(0)  # Reset file pointer
        try:
            file_content = file.read()
            images = convert_from_bytes(file_content)
            logger.info(f"Converted page {page_num + 1} to {len(images)} images")
        except Exception as e:
            logger.error(f"Error converting page {page_num + 1} to image: {str(e)}")
            images = []
        
        combined_image_text = ""
        for img_idx, img in enumerate(images):
            try:
                # Save the image with a unique filename
                local_file_storage_path=os.getenv("LOCAL_FILE_STORAGE_PATH") or ""
                image_filename = f"{local_file_storage_path}pdf_images/{document_id}/{uuid.uuid4().hex}.png"
                dir_path = os.path.dirname(default_file_storage.path(image_filename))
                os.makedirs(dir_path, exist_ok=True)
                
                img_byte_arr = io.BytesIO()
                img.save(img_byte_arr, format='PNG')
                img_byte_arr.seek(0)
                
                # Validate and save the image
                img_bytes = img_byte_arr.getvalue()
                if not is_valid_image(img_bytes):
                    logger.warning(f"Skipping invalid image {img_idx + 1} on page {page_num + 1}")
                    continue
                default_file_storage.save(image_filename, img_byte_arr)
                logger.info(f"Saved image {img_idx + 1} on page {page_num + 1} as {image_filename}")
                
                # Analyze the image with LLM (single, efficient processing step)
                img_byte_arr.seek(0)
                analysis_result = image_analyzer.analyze_image(img_byte_arr.getvalue())
                logger.info(f"Analyzed image {img_idx + 1} on page {page_num + 1}")
                
                # Verify the structure of the analysis result
                if not isinstance(analysis_result, dict) or "text_content" not in analysis_result or "description" not in analysis_result:
                    logger.error(f"Invalid analysis result for image {img_idx + 1} on page {page_num + 1}")
                    continue
                
                # Create an image record using the analysis data
                image = Image(
                    path=image_filename,
                    caption=f"Page {page_num + 1}, Image {img_idx + 1}",
                    text_snippets=analysis_result["text_content"],
                    description=analysis_result["description"],
                    source_document_id=document_id,
                )
                images_list.append(image)
                combined_image_text += analysis_result["text_content"] + "\n"
            except Exception as e:
                logger.error(f"Error processing image {img_idx + 1} on page {page_num + 1}: {str(e)}")
        
        # Combine text sources for the page
        if embedded_text and combined_image_text.strip():
            page_text = embedded_text if len(embedded_text) > len(combined_image_text) else combined_image_text.strip()
            logger.info(f"Using {'embedded' if len(embedded_text) > len(combined_image_text) else 'LLM'} text for page {page_num + 1}")
        elif embedded_text:
            page_text = embedded_text
            logger.info(f"Using only embedded text for page {page_num + 1}")
        elif combined_image_text.strip():
            page_text = combined_image_text.strip()
            logger.info(f"Using only LLM text for page {page_num + 1}")
        else:
            page_text = ""
            logger.warning(f"No text extracted from page {page_num + 1}")
        
        text_content.append(page_text)
    
    logger.info(f"Completed PDF extraction for document ID: {document_id}")
    return "\n\n".join(text_content), images_list






def extract_text_from_docx(file: IO) -> str:
    document = docx.Document(file)
    full_text = []
    for paragraph in document.paragraphs:
        full_text.append(paragraph.text)
    return "\n\n".join(full_text)


def extract_text_from_pptx(file: IO) -> str:
    presentation = pptx.Presentation(file)
    full_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n\n".join(full_text)


def extract_text_from_xlsx(file: IO) -> str:
    wb = openpyxl.load_workbook(file)
    full_text = []
    for sheet in wb.worksheets:
        full_text.append(f"Sheet: {sheet.title}")
        sheet_string = "\n".join(
            ",".join(map(str, row))
            for row in sheet.iter_rows(values_only=True)
        )
        full_text.append(sheet_string)
    return "\n\n".join(full_text)


def img_to_bytes(img) -> bytes:
    """Convert PIL Image to bytes."""
    img_byte_arr = io.BytesIO()
    img.save(img_byte_arr, format='PNG')
    img_byte_arr.seek(0)
    logger.debug(f"Converted image to bytes, size: {len(img_byte_arr.getvalue())} bytes")
    return img_byte_arr.getvalue()


# Add a new function to check if an image is valid
def is_valid_image(img_bytes: bytes) -> bool:
    """Check if the image bytes represent a valid image."""
    try:
        img = PILImage.open(io.BytesIO(img_bytes))
        img.verify()  # Verify the image data
        logger.debug(f"Image validation successful, size: {img.size}")
        return True
    except Exception as e:
        logger.error(f"Image validation failed: {str(e)}")
        return False
